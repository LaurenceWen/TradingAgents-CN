from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _read_version(repo_root: Path) -> str:
    try:
        return (repo_root / "VERSION").read_text(encoding="utf-8").strip()
    except Exception:
        return "unknown"


def _set_text_run(run, font_name: str, font_size: int, rgb: tuple[int, int, int], bold: bool = False):
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*rgb)


def _add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.8), Inches(1.4))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    _set_text_run(r, "Microsoft YaHei", 40, (31, 78, 121), bold=True)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.7), Inches(11.5), Inches(1.0))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.LEFT
        sr = sp.add_run()
        sr.text = subtitle
        _set_text_run(sr, "Microsoft YaHei", 20, (64, 64, 64), bold=False)


def _add_section_header(slide, header: str):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = header
    _set_text_run(r, "Microsoft YaHei", 30, (31, 78, 121), bold=True)


def _add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str]):
    """Add bullets with a compact, readable layout.

    NOTE: We intentionally prefix lines with a bullet symbol to improve visual scanning,
    and auto-adjust font size for longer lists.
    """

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    # Tighter, more consistent margins (default margins vary by layout)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.clear()

    # Auto font sizing by number of lines
    n = max(len(bullets), 1)
    font_size = 20
    if n >= 6:
        font_size = 18
    if n >= 8:
        font_size = 16

    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text = (line or "").strip()
        if text and not text.startswith("•"):
            text = f"• {text}"
        p.text = text
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.15
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(40, 40, 40)


def _add_two_columns(slide, header: str, left_bullets: list[str], right_bullets: list[str]):
    _add_section_header(slide, header)
    _add_bullets(slide, 0.9, 1.4, 5.9, 5.6, left_bullets)
    _add_bullets(slide, 7.0, 1.4, 5.9, 5.6, right_bullets)


def _add_architecture(slide):
    _add_section_header(slide, "产品形态与交付（企业级）")

    def box(x, y, w, h, text, fill):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
        shape.line.color.rgb = RGBColor(200, 200, 200)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        _set_text_run(r, "Microsoft YaHei", 16, (255, 255, 255), bold=True)
        return shape

    def arrow(x, y, w, h):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
        shape.line.color.rgb = RGBColor(220, 220, 220)
        return shape

    box(0.7, 1.55, 2.2, 0.9, "用户/浏览器", (91, 155, 213))
    arrow(3.0, 1.78, 0.6, 0.45)
    box(3.7, 1.55, 3.1, 0.9, "前端：Vue 3 + Element Plus", (47, 85, 151))
    arrow(6.9, 1.78, 0.6, 0.45)
    box(7.6, 1.55, 3.0, 0.9, "后端：FastAPI API", (31, 78, 121))

    arrow(10.7, 1.78, 0.6, 0.45)
    box(11.4, 1.55, 1.6, 0.9, "Worker", (31, 78, 121))

    box(3.7, 3.0, 4.0, 0.85, "实时推送：SSE / WebSocket", (112, 173, 71))
    box(8.0, 3.0, 5.0, 0.85, "任务：分析队列 / 批量处理", (112, 173, 71))

    box(3.7, 4.35, 3.8, 0.9, "Redis：队列 + 缓存", (237, 125, 49))
    box(7.7, 4.35, 3.8, 0.9, "MongoDB：数据存储", (237, 125, 49))
    box(11.7, 4.35, 1.3, 0.9, "统一分析引擎\n+ 工作流", (91, 155, 213))

    note = slide.shapes.add_textbox(Inches(0.9), Inches(6.05), Inches(12.4), Inches(0.95))
    tf = note.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "交付：Docker Compose 私有化一键部署；权限/审计日志/任务中心/实时推送，满足机构落地与运营。"
    _set_text_run(r, "Microsoft YaHei", 16, (70, 70, 70), bold=False)


def _add_why_now_and_pain(slide):
    _add_section_header(slide, "Why now：大模型进场，但金融落地仍缺‘产品化’")
    _add_bullets(
        slide,
        0.95,
        1.35,
        12.0,
        5.8,
        [
            "大模型能力已足够强，但‘投研→决策→复盘’是流程密集型工作，靠聊天很难规模化",
            "落地关键不在模型本身，而在流程、追溯、管控、交付：能否形成可复制的方法论",
            "机构真实痛点：口径不统一、结果不可复现、成本不可控、合规与审计难",
            "机会：把投研方法论做成工作流与模板资产，沉淀组织能力与数据资产",
        ],
    )


def _add_solution_overview(slide):
    _add_section_header(slide, "我们的解法：把投研方法论产品化（工作流 + 模板 + 任务中心）")

    def box(x, y, w, h, text, fill, font_size=16):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
        shape.line.color.rgb = RGBColor(210, 210, 210)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        _set_text_run(r, "Microsoft YaHei", font_size, (255, 255, 255), bold=True)
        return shape

    def arrow(x, y, w, h):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(225, 225, 225)
        shape.line.color.rgb = RGBColor(225, 225, 225)
        return shape

    # Row 1: Entry points
    box(0.8, 1.55, 3.5, 0.9, "入口（业务闭环）\n单股 / 批量 / 持仓 / 复盘", (91, 155, 213), font_size=15)
    arrow(4.45, 1.78, 0.55, 0.45)
    box(5.1, 1.55, 3.9, 0.9, "中枢（方法论引擎）\n统一分析引擎 + 多智能体工作流", (31, 78, 121), font_size=15)
    arrow(9.15, 1.78, 0.55, 0.45)
    box(9.8, 1.55, 3.2, 0.9, "输出\n报告 / 证据链 / 风险提示", (47, 85, 151), font_size=15)

    # Row 2: Operability / governance
    box(0.8, 3.05, 4.0, 0.85, "可运营：模板体系\n版本/灰度/回滚/偏好", (112, 173, 71), font_size=15)
    box(5.1, 3.05, 3.9, 0.85, "可追溯：任务中心\n进度/耗时/失败原因/重跑", (112, 173, 71), font_size=15)
    box(9.2, 3.05, 3.8, 0.85, "可管控：多模型治理\n选择/参数/用量/成本", (112, 173, 71), font_size=15)

    box(0.8, 4.45, 12.2, 0.85, "可交付：私有化部署（Docker Compose） + 权限/审计日志", (237, 125, 49), font_size=16)

    note = slide.shapes.add_textbox(Inches(0.95), Inches(5.65), Inches(12.1), Inches(1.1))
    tf = note.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "一句话：不是‘聊天看盘’，而是把投研流程做成可配置、可追溯、可运营、可交付的企业级产品。"
    _set_text_run(r, "Microsoft YaHei", 16, (70, 70, 70), bold=False)


def _add_agents(slide, repo_root: Path):
    _add_section_header(slide, "多智能体协作（示意）")
    bullets = [
        "模拟“投研团队”协作：分析师 → 研究员（多观点对辩）→ 风险审查 → 交易/复盘",
        "多维度输入：基本面 / 技术面 / 新闻与舆情 / 多源行情数据",
        "输出：结构化结论、证据链、风险提示、报告与复盘案例沉淀",
    ]
    _add_bullets(slide, 0.9, 1.35, 7.0, 2.5, bullets)

    assets = repo_root / "assets"
    images = [
        ("分析师", assets / "analyst.png"),
        ("研究员", assets / "researcher.png"),
        ("风控", assets / "risk.png"),
        ("交易员", assets / "trader.png"),
    ]
    x0 = 0.95
    y0 = 4.0
    w = 2.7
    h = 2.6
    for i, (label, path) in enumerate(images):
        x = x0 + i * (w + 0.25)
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(x + 0.65), Inches(y0 + 0.15), width=Inches(1.4))
        tbox = slide.shapes.add_textbox(Inches(x), Inches(y0 + 1.7), Inches(w), Inches(0.6))
        tf = tbox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        _set_text_run(r, "Microsoft YaHei", 18, (40, 40, 40), bold=True)


def _add_v2_capabilities(slide):
    _add_two_columns(
        slide,
        "核心能力（v2.0）",
        left_bullets=[
            "统一分析引擎：跨入口字段对齐，结果一致、可复用",
            "可视化工作流：节点/边编排，支持条件与循环（对辩/分支）",
            "任务中心：进度、节点输出、耗时与重跑（可追溯、可审计）",
        ],
        right_bullets=[
            "模板体系：系统模板 + 用户模板，版本/灰度/回滚",
            "多模型治理：统一接入，参数/用量/成本可控（供应商可替换）",
            "闭环：持仓分析 + 交易复盘 + 案例库沉淀",
        ],
    )


def _add_workflow_customization(slide):
    _add_section_header(slide, "流程定制：可视化工作流 + 统一分析引擎")
    _add_bullets(
        slide,
        0.95,
        1.35,
        12.0,
        5.8,
        [
            "把投研链路抽象为‘节点+边’：可视化增删节点、调整顺序",
            "每个节点可配置模型/参数：能力与成本按需组合",
            "支持分支/循环 + 模板复用：保存/版本化/一键执行，沉淀方法论资产",
        ],
    )


def _add_prompt_customization(slide):
    _add_section_header(slide, "提示词定制：模板体系（可运营、可控、可回滚）")

    # Left: concise bullets
    _add_bullets(
        slide,
        0.95,
        1.35,
        6.3,
        5.8,
        [
            "系统模板：覆盖分析/研究/风控/交易/复盘等 Agent 类别",
            "用户模板：派生/启停/版本管理，统一团队口径",
            "与工作流联动：流程 × 模板 = 方法论可交付",
        ],
    )

    # Right: visual operability chips (less text, stronger layout)
    caption = slide.shapes.add_textbox(Inches(7.45), Inches(1.35), Inches(5.5), Inches(0.5))
    ctf = caption.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.LEFT
    cr = cp.add_run()
    cr.text = "运营机制（团队可复用）"
    _set_text_run(cr, "Microsoft YaHei", 22, (31, 78, 121), bold=True)

    def chip(x, y, text, fill=(112, 173, 71)):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(2.55),
            Inches(0.7),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
        shape.line.color.rgb = RGBColor(255, 255, 255)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        _set_text_run(r, "Microsoft YaHei", 18, (255, 255, 255), bold=True)

    x0, y0 = 7.45, 2.05
    chip(x0, y0, "版本管理")
    chip(x0 + 2.75, y0, "灰度发布")
    chip(x0, y0 + 0.9, "一键回滚")
    chip(x0 + 2.75, y0 + 0.9, "分层权限")

    # Small note to connect to governance value
    note = slide.shapes.add_textbox(Inches(7.45), Inches(4.05), Inches(5.5), Inches(1.2))
    ntf = note.text_frame
    ntf.clear()
    np = ntf.paragraphs[0]
    np.alignment = PP_ALIGN.LEFT
    nr = np.add_run()
    nr.text = "让提示词从‘个人手艺’变成‘组织资产’：可追溯、可运营、可复制。"
    _set_text_run(nr, "Microsoft YaHei", 18, (70, 70, 70), bold=False)


def _add_closed_loop(slide):
    _add_two_columns(
        slide,
        "闭环能力：持仓分析 + 交易复盘（PRO）",
        left_bullets=[
            "组合体检：市值/成本/盈亏、行业分布、集中度与风险等级",
            "持仓流水：加减仓、分红、拆并股、成本调整等完整记录",
            "单只持仓诊断：结合买入成本/权重/持仓天数给出建议",
        ],
        right_bullets=[
            "单笔复盘：纪律、情绪、仓位、执行偏差与改进建议",
            "阶段复盘：周/月/季度汇总胜率、回撤、盈亏比，生成行动计划",
            "案例库沉淀：复盘可保存打标签，形成组织/个人知识资产",
        ],
    )


def _add_why_us(slide):
    _add_section_header(slide, "为什么是我们（对比通用大模型）")
    _add_bullets(
        slide,
        0.95,
        1.35,
        12.0,
        5.8,
        [
            "通用大模型：强在对话；弱在流程治理（追溯/管控/交付）",
            "我们：工作流驱动——把‘分析→决策→复盘’沉淀为可复用方法论",
            "我们：任务中心节点级留痕——审计/复现/重跑/迭代都更容易",
            "我们：多模型治理——统一接入，供应商可替换，成本可控",
        ],
    )


def _add_portfolio(slide):
    _add_section_header(slide, "持仓分析：从“看股票”走向“看组合”（PRO）")
    _add_bullets(
        slide,
        0.95,
        1.35,
        12.0,
        5.8,
        [
            "持仓管理：新增/更新/删除、批量导入，支持真实/模拟数据源",
            "持仓操作：加仓、减仓、分红、拆股/合股、调整成本等完整流水",
            "组合体检：总市值/成本/盈亏、行业分布、集中度与风险等级",
            "单只持仓诊断：结合买入成本、持仓天数与权重给出操作建议",
            "统一任务中心：分析过程异步执行，可查询进度并留存历史",
        ],
    )


def _add_review(slide):
    _add_section_header(slide, "交易复盘：操作复盘 + 阶段复盘（PRO）")
    _add_bullets(
        slide,
        0.95,
        1.35,
        12.0,
        5.8,
        [
            "单笔/完整交易复盘：从交易记录抽取全链路，输出纪律、情绪、仓位等点评",
            "阶段复盘：按周/月/季度汇总胜率、盈亏比、回撤、费用等，生成行动计划",
            "案例库沉淀：复盘结果可保存为案例并打标签，形成个人教科书",
            "交易计划联动：复盘可关联交易计划，评估“是否按规则执行”",
            "统一工作流驱动：复盘分析同样由多智能体工作流执行",
        ],
    )


def _add_investor_highlights(slide):
    _add_section_header(slide, "投资人真正买单的亮点（本项目对应）")
    _add_bullets(
        slide,
        0.95,
        1.35,
        12.0,
        5.8,
        [
            "不是“又一个ChatGPT看盘”：把投研方法论产品化（工作流 + 模板 + 任务中心），可复制、可交付",
            "不是“只会给结论”：过程可追溯（每个节点输出可见），便于合规、审计、复现与迭代",
            "不是“只分析个股”：闭环覆盖持仓→操作→复盘→案例库，沉淀用户数据与使用粘性",
            "不是“绑死某一家模型”：多LLM统一接入 + 成本与用量可控，机构落地更稳",
            "不是“演示品”：前后端 + 队列/缓存/推送 + 权限/日志，具备可运营与规模化基础",
        ],
    )


def _add_moat(slide):
    _add_two_columns(
        slide,
        "护城河（为什么你能持续赢）",
        left_bullets=[
            "工作流资产：模板/节点/参数沉淀为可复用方法论库",
            "模板资产：版本化、灰度、回滚，形成组织级标准口径",
            "复盘与案例库：行为与结果沉淀为私有数据资产，越用越强",
        ],
        right_bullets=[
            "合规友好交付：私有化部署、权限与审计日志，机构采用门槛低",
            "多模型可替换：供应商变化不伤筋骨，成本优化与议价空间大",
            "可扩展：新增节点/工具/Agent 成本低，迭代速度快",
        ],
    )


def _add_business_model(slide):
    _add_two_columns(
        slide,
        "商业化与合作",
        left_bullets=[
            "ToB：私有化部署 / 授权订阅 / 定制开发 / 集成服务",
            "交付物：工作流模板、模板体系、报告导出与数据对接",
            "增值点：数据治理、成本优化、合规模板与审计能力",
        ],
        right_bullets=[
            "ToC：Pro 订阅（批量分析、持仓/复盘、报告与历史沉淀）",
            "教育版：课程/训练营/案例库与作业评测（可选）",
            "合作方：券商投研、资管、量化团队、教育培训机构",
        ],
    )


def _add_roadmap_cta(slide):
    _add_section_header(slide, "路线图与合作方式")
    _add_bullets(
        slide,
        0.95,
        1.35,
        12.0,
        5.8,
        [
            "近期：更多行业工作流模板与报告标准化（提升交付效率）",
            "中期：组织协作（团队模板共享/权限分层/审计增强）",
            "远期：模板市场与生态扩展（更多资产类别与场景）",
            "合作方式：POC试点 → 私有化部署/授权订阅 → 深度定制与集成",
        ],
    )

    footer = slide.shapes.add_textbox(Inches(0.95), Inches(6.45), Inches(12.0), Inches(0.6))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "谢谢｜欢迎交流：投资/合作/私有化落地/联合交付"
    _set_text_run(r, "Microsoft YaHei", 18, (31, 78, 121), bold=True)


def build_ppt(repo_root: Path, output_path: Path) -> Path:
    version = _read_version(repo_root)
    today = datetime.now().strftime("%Y-%m-%d")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(
        slide,
        "TradingAgents-CN Pro 项目介绍",
        f"企业级投研工作流产品（多智能体 + 大模型） | 版本 {version} | {today}",
    )
    cover_note = slide.shapes.add_textbox(Inches(0.85), Inches(4.35), Inches(11.8), Inches(1.7))
    tf = cover_note.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "定位：投研与复盘流程的企业级工作流产品；合规友好，不提供实盘交易指令。"
    _set_text_run(r, "Microsoft YaHei", 20, (55, 55, 55), bold=False)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_two_columns(
        slide,
        "价值主张（我们交付什么价值）",
        left_bullets=[
            "标准化：把投研方法论沉淀为可复用工作流模板",
            "可追溯：任务中心留痕到节点级（输出/耗时/失败原因）",
            "可管控：多模型统一接入，参数/用量/成本可控",
        ],
        right_bullets=[
            "可运营：模板版本化（灰度/回滚），统一口径与偏好",
            "可交付：Docker Compose 私有化部署 + 权限/审计日志",
            "可闭环：持仓分析 + 交易复盘 + 案例库沉淀",
        ],
    )

    # 3) Why now & pain
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_why_now_and_pain(slide)

    # 4) Solution overview (the storyline anchor)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_solution_overview(slide)

    # 5) Product & delivery
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_architecture(slide)

    # 6) v2 capabilities overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_v2_capabilities(slide)

    # 7) Workflow customization
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_workflow_customization(slide)

    # 8) Prompt template operability
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_prompt_customization(slide)

    # 9) Closed-loop capabilities
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_closed_loop(slide)

    # 10) Why us
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_why_us(slide)

    # 11) Moat
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_moat(slide)

    # 12) Business model
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_business_model(slide)

    # 13) Roadmap & CTA
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_roadmap_cta(slide)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main():
    repo_root = Path(__file__).resolve().parents[2]
    version = _read_version(repo_root)
    now = datetime.now()
    today = now.strftime("%Y-%m-%d")
    time_tag = now.strftime("%H%M%S")
    out = repo_root / "copyright_submission" / f"TradingAgentsCN_项目介绍_{version}_{today}_{time_tag}.pptx"
    build_ppt(repo_root=repo_root, output_path=out)
    print(str(out))


if __name__ == "__main__":
    main()
